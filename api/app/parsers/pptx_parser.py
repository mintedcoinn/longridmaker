from __future__ import annotations

import mimetypes
import tempfile
from io import BytesIO
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.table import Table
from PIL import Image

_CAPTION_GAP_THRESHOLD = 1_200_000
_CAPTION_HORIZONTAL_OVERLAP_RATIO = 0.35


@dataclass
class ShapeBox:
    text: str
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height


def parse_pptx(file_path: Path, assets_dir: Path | None = None) -> dict:
    """Parse a PPTX file and return structured JSON slide-by-slide."""
    prs = Presentation(str(file_path))
    resolved_assets_dir = _resolve_assets_dir(assets_dir)
    slides = []
    image_number = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name if slide.slide_layout else None
        notes_text = _extract_notes_text(slide)

        shape_records = []
        text_boxes: list[ShapeBox] = []
        tables = []
        images = []

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            shape_data = _extract_shape(shape, shape_idx)
            if shape_data is None:
                continue

            shape_records.append(shape_data)

            if shape_data["type"] == "text":
                text_boxes.append(
                    ShapeBox(
                        text=shape_data["text"],
                        left=shape_data["left"],
                        top=shape_data["top"],
                        width=shape_data["width"],
                        height=shape_data["height"],
                    )
                )
            elif shape_data["type"] == "table":
                tables.append(shape_data)
            elif shape_data["type"] == "image":
                image_number += 1
                caption = _find_caption_for_image(shape_data, text_boxes)
                image_path = _save_image(shape, resolved_assets_dir, slide_idx, image_number)
                shape_data["image_number"] = image_number
                shape_data["image_path"] = str(image_path) if image_path else None
                shape_data["caption"] = caption
                images.append(shape_data)

        slides.append(
            {
                "slide_number": slide_idx,
                "layout": {
                    "name": layout_name,
                },
                "notes": {
                    "text": notes_text,
                },
                "shapes": shape_records,
                "images": images,
                "tables": tables,
            }
        )

    return {
        "file_type": "pptx",
        "slides": slides,
        "metadata": {
            "total_slides": len(slides),
            "total_images": image_number,
            "assets_dir": str(resolved_assets_dir),
        },
    }


def _resolve_assets_dir(assets_dir: Path | None) -> Path:
    if assets_dir is not None:
        assets_dir.mkdir(parents=True, exist_ok=True)
        return assets_dir

    return Path(tempfile.mkdtemp(prefix="file-parser-pptx-assets-"))


def _extract_shape(shape: BaseShape, shape_index: int) -> dict | None:
    """Extract content from a single shape."""
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)

    if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}:
        image = shape.image
        descr = _get_picture_description(shape)
        return {
            "type": "image",
            "shape_index": shape_index,
            "name": shape.name,
            "caption_hint": descr,
            "content_type": image.content_type,
            "filename": image.filename,
            "extension": image.ext,
            "sha1": image.sha1,
            "size": list(image.size),
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        }

    if shape.has_text_frame:
        text = _extract_text(shape)
        if text:
            return {
                "type": "text",
                "shape_index": shape_index,
                "name": shape.name,
                "text": text,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
            }

    if shape.has_table:
        table: Table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        return {
            "type": "table",
            "shape_index": shape_index,
            "name": shape.name,
            "rows": rows,
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        }

    return None


def _extract_text(shape: BaseShape) -> str:
    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_notes_text(slide) -> str:
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""

    text_frame = getattr(notes_slide, "notes_text_frame", None)
    if text_frame is None:
        return ""

    paragraphs = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    return "\n".join(paragraphs)


def _find_caption_for_image(image_shape: dict, text_boxes: list[ShapeBox]) -> str | None:
    if not text_boxes:
        return image_shape.get("caption_hint")

    image_left = image_shape["left"]
    image_top = image_shape["top"]
    image_right = image_left + image_shape["width"]
    image_bottom = image_top + image_shape["height"]
    image_width = max(image_shape["width"], 1)
    image_center_x = image_left + image_width / 2

    candidates = []
    for text_box in text_boxes:
        horizontal_overlap = _horizontal_overlap(image_left, image_right, text_box.left, text_box.right)
        overlap_ratio = horizontal_overlap / max(min(image_width, text_box.width), 1)
        is_below = text_box.top >= image_bottom
        gap = text_box.top - image_bottom
        is_close = gap <= _CAPTION_GAP_THRESHOLD

        if is_below and is_close and overlap_ratio >= _CAPTION_HORIZONTAL_OVERLAP_RATIO:
            distance = gap + abs((text_box.left + text_box.width / 2) - image_center_x) / 5
            candidates.append((distance, text_box.text))

    if candidates:
        candidates.sort(key=lambda item: item[0])
        return candidates[0][1]

    caption_hint = image_shape.get("caption_hint")
    if caption_hint:
        return caption_hint

    return None


def _horizontal_overlap(left_a: int, right_a: int, left_b: int, right_b: int) -> int:
    return max(0, min(right_a, right_b) - max(left_a, left_b))


def _save_image(shape: BaseShape, assets_dir: Path, slide_number: int, image_number: int) -> Path | None:
    try:
        image = shape.image
    except Exception:
        return None

    suffix = f".{image.ext}" if image.ext else Path(image.filename or "").suffix
    if not suffix:
        suffix = mimetypes.guess_extension(image.content_type) or ".bin"

    normalized_suffix = suffix.lower()
    image_path = assets_dir / f"pptx_slide_{slide_number:04d}_image_{image_number:04d}"

    if normalized_suffix in {".wmf", ".emf"}:
        converted = _save_preview_png(image.blob, image_path.with_suffix(".png"))
        if converted is not None:
            return converted

    image_path = image_path.with_suffix(suffix)
    image_path.write_bytes(image.blob)
    return image_path


def _save_preview_png(blob: bytes, target_path: Path) -> Path | None:
    try:
        with Image.open(BytesIO(blob)) as image:
            image.load()
            image.save(target_path, format="PNG")
        return target_path
    except Exception:
        return None


def _get_picture_description(shape: BaseShape) -> str | None:
    try:
        descr = shape._element.nvPicPr.cNvPr.get("descr")
        if descr:
            return descr.strip()
    except Exception:
        return None
    return None
